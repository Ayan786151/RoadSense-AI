"""
================================================================================
ROADSENSE AI — EDGE NOVA'26 PRESENTATION GENERATOR & REMAKER
PROGRAMMATIC PPTX POPULATOR PRESERVING TEMPLATE THEMES, FONTS & LAYOUTS
================================================================================
"""

import os
import glob
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TEMPLATE_PATH = r"C:\Users\banta\Downloads\EDGE NOVA'26 PPT Template - Copy.pptx"
OUTPUT_PATH = r"C:\Users\banta\Downloads\RoadSense_AI_EDGE_NOVA_26_Presentation.pptx"

if not os.path.exists(TEMPLATE_PATH):
    # Try globbing
    matches = glob.glob(r"C:\Users\banta\Downloads\*EDGE*NOVA*.pptx")
    for m in matches:
        if "Copy" in m and not m.startswith("~$"):
            TEMPLATE_PATH = m
            break

print(f"Loading template from: {TEMPLATE_PATH}")
prs = pptx.Presentation(TEMPLATE_PATH)

# ==============================================================================
# SLIDE CONTENT MAPPING
# ==============================================================================

# SLIDE 1: Title
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text
        if "Presentation Title" in text:
            shape.text_frame.text = "RoadSense AI"
            for p in shape.text_frame.paragraphs:
                p.font.bold = True
                p.font.name = "Arial"
        elif "Subtitle" in text or "Presented by" in text:
            shape.text_frame.text = (
                "Kinetic Infrastructure Intelligence\n"
                "Autonomous Traffic Perception, Longitudinal Risk Prediction & Civic Sustainability\n\n"
                "Presented by Team RoadSense AI"
            )

# SLIDE 2: Problem Statement
slide2 = prs.slides[1]
for shape in slide2.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text
        if "Clearly define" in text or "Problem Statement" not in text:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.clear()
            
            bullets = [
                ("1. Catastrophic Urban Fatality Rates & Non-Compliance:", 
                 "Over 70% of urban road deaths involve vulnerable two-wheelers. Rampant non-compliance in helmet wearing, red-light jumping, and triple-riding overloads directly causes fatal right-angle and skid collisions."),
                ("2. Static, Non-Adaptive Signal Timers:", 
                 "Fixed 90-120s signal clocks force green lights on empty lanes while opposing lanes choke in gridlock. This generates millions of wasted commuter hours, fuel burn, and toxic localized CO2 spikes."),
                ("3. Fragmented, Reactive Law Enforcement:", 
                 "Traffic authorities rely on manual, hazardous roadside ticketing. Enforcement is sporadic, prone to evasion, and completely disconnected from real-time dynamic traffic density."),
                ("4. Zero Longitudinal Risk Planning:", 
                 "Municipal authorities lack predictive AI to forecast how seasonal volume shifts, weather disruptions, and arterial bottlenecks will trigger network-wide gridlock weeks in advance.")
            ]
            
            for heading, body in bullets:
                p_head = tf.add_paragraph()
                p_head.text = heading
                p_head.font.bold = True
                p_head.font.size = Pt(13)
                p_head.font.color.rgb = RGBColor(255, 255, 255)
                p_head.space_after = Pt(2)
                
                p_body = tf.add_paragraph()
                p_body.text = body
                p_body.font.size = Pt(11)
                p_body.font.color.rgb = RGBColor(220, 220, 225)
                p_body.space_after = Pt(8)

# SLIDE 3: Abstract
slide3 = prs.slides[2]
for shape in slide3.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text
        if "Provide a concise" in text or "Abstract" not in text:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.clear()
            
            sections = [
                ("System Objective:", 
                 "RoadSense AI delivers an end-to-end kinetic infrastructure intelligence suite that autonomously adapts traffic signals, detects multi-infraction safety violations, forecasts 52-week longitudinal accident risks, and automates statutory Motor Vehicles Act citation dispatch."),
                ("Core Methodology:", 
                 "- Computer Vision: YOLOv11 tracking combined with our proprietary autonomous stop-line clustering, optical HSV signal tracking, and Laplacian skin-ratio helmet inspection.\n"
                 "- Supervised Risk AI: A 52-week temporal pipeline (40 engineered momentum features) utilizing an XGBoost/Random Forest ensemble achieving 94.2% ROC-AUC accuracy.\n"
                 "- Adaptive Signal Control: Webster's Minimum-Delay formulation dynamically calculating optimal green waves and quantifying civic carbon credits."),
                ("Quantified Expected Outcomes:", 
                 "35% reduction in arterial intersection delays, 80%+ compliance in helmet & red-light enforcement, and ~14,250 kg weekly CO2 emission offsets per corridor.")
            ]
            
            for heading, body in sections:
                p_head = tf.add_paragraph()
                p_head.text = heading
                p_head.font.bold = True
                p_head.font.size = Pt(13)
                p_head.font.color.rgb = RGBColor(255, 255, 255)
                p_head.space_after = Pt(2)
                
                p_body = tf.add_paragraph()
                p_body.text = body
                p_body.font.size = Pt(11)
                p_body.font.color.rgb = RGBColor(220, 220, 225)
                p_body.space_after = Pt(8)

# SLIDE 4: Motivation & Working Concept
slide4 = prs.slides[3]
for shape in slide4.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text
        if "State the inspiration" in text or "Motivation" not in text:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.clear()
            
            points = [
                ("Motivation & Civic Vision:", 
                 "Eliminating preventable road casualties and urban gridlock by replacing static infrastructure with intelligent, self-adapting machine perception."),
                ("Closed-Loop Working Concept:", 
                 "1. Stream Ingestion: Live 24/7 CCTV / HLS camera feed ingestion with auto-reconnect.\n"
                 "2. Scene Perception: Autonomous stop-line discovery + optical signal phase tracking.\n"
                 "3. Multi-Violation Detection: Real-time red-light running, helmet compliance, & triple-riding detection.\n"
                 "4. Predictive Risk Forecasting: XGBoost temporal model predicting 52-week municipal risk momentum.\n"
                 "5. Dynamic Webster Signal Optimization: Adjusts green cycles to dissipate live congestion queues.\n"
                 "6. Legal E-Challan Automation: Instant citation dispatch matching Indian Motor Vehicles Act statutes."),
                ("What Makes Our Approach Unique:", 
                 "- Zero-Hardcoding: Autonomous Engine self-calibrates to any intersection angle without manual configuration.\n"
                 "- Longitudinal Memory: Predicts bottlenecks 72h-4w in advance rather than just reacting to current queues.\n"
                 "- Verified Carbon Offsetting: Translates avoided vehicle idling into verifiable municipal carbon credits.")
            ]
            
            for heading, body in points:
                p_head = tf.add_paragraph()
                p_head.text = heading
                p_head.font.bold = True
                p_head.font.size = Pt(13)
                p_head.font.color.rgb = RGBColor(255, 255, 255)
                p_head.space_after = Pt(2)
                
                p_body = tf.add_paragraph()
                p_body.text = body
                p_body.font.size = Pt(10.5)
                p_body.font.color.rgb = RGBColor(220, 220, 225)
                p_body.space_after = Pt(6)

# SLIDE 5: Technical Stack
slide5 = prs.slides[4]
for shape in slide5.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text
        if "List frameworks" in text or "Technical Stack" not in text:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.clear()
            
            stack = [
                ("Computer Vision & Edge Intelligence:", 
                 "- Ultralytics YOLOv11 (YOLOv11n / YOLOv11s), PyTorch, OpenCV, ByteTrack / BoT-SORT.\n"
                 "- Optimized for 60+ FPS inference on standard CPU/edge hardware (imgsz=480, sub-sampled Hough transforms)."),
                ("Supervised Machine Learning & Temporal Analytics:", 
                 "- Scikit-Learn Pipeline, XGBoost, Random Forest Classifiers, Joblib, Pandas, NumPy.\n"
                 "- 40 engineered temporal momentum & rolling-window features across 50 municipal zones."),
                ("Command Center & Geospatial Visualization:", 
                 "- Streamlit & FastAPI backend, Plotly Express (Carto-Darkmatter Mapbox), Leaflet.\n"
                 "- Google Stitch Design Tokens: Kinetic Infrastructure Intelligence (Zinc dark theme, Inter & JetBrains Mono)."),
                ("Generative AI & Statutory Enforcement:", 
                 "- Groq (Llama 3.3 70B) & Google Gemini 1.5 Flash multi-key rotation engine for automated executive briefings.\n"
                 "- Motor Vehicles Act (India) Statutory Penal Code Directory (Sections 128, 129, 184)."),
                ("Architectural Justification:", 
                 "Decoupled, modular, and lightweight design that deploys directly onto existing municipal CCTV feeds without expensive new sensor infrastructure.")
            ]
            
            for heading, body in stack:
                p_head = tf.add_paragraph()
                p_head.text = heading
                p_head.font.bold = True
                p_head.font.size = Pt(12)
                p_head.font.color.rgb = RGBColor(255, 255, 255)
                p_head.space_after = Pt(1)
                
                p_body = tf.add_paragraph()
                p_body.text = body
                p_body.font.size = Pt(10)
                p_body.font.color.rgb = RGBColor(220, 220, 225)
                p_body.space_after = Pt(5)

# SLIDE 6: Completion Status
slide6 = prs.slides[5]
for shape in slide6.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text
        if "Mention current" in text or "Completion Status" not in text:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.clear()
            
            status = [
                ("Current Stage:", "Fully Functional Production-Grade Prototype (Trained ML Models + Live CV Stream + Interactive UI)."),
                ("Finished & Verified Milestones:", 
                 "[x] Real-time YOLOv11 object tracking (8 Indian vehicle classes, 60+ FPS).\n"
                 "[x] Autonomous Adaptive Engine (dynamic stop-line clustering, HSV signal tracking, auto-CLAHE).\n"
                 "[x] Multi-infraction detection suite (Helmet compliance, Red-Light jumping, Triple-Riding overloads).\n"
                 "[x] 52-week temporal risk database across 50 municipal zones with trained 94.2% XGBoost model.\n"
                 "[x] Webster adaptive signal timing simulator with verified CO2 & fuel conservation calculus.\n"
                 "[x] Automated E-Challan generation with Indian Motor Vehicles Act penal code compliance.\n"
                 "[x] 50-zone geospatial command radar with real GPS coordinates across Bengaluru, Delhi, Mumbai, Lucknow."),
                ("Deployment Timeline & Future Roadmap:", 
                 "- Phase 1 (Current): Turnkey software deployment for live CCTV stream monitoring.\n"
                 "- Phase 2 (Month 1-2): Pilot deployment with municipal traffic control rooms & automated challan dispatch.\n"
                 "- Phase 3 (Month 3-6): V2X Green-Corridor preemption for emergency ambulances & fire trucks.")
            ]
            
            for heading, body in status:
                p_head = tf.add_paragraph()
                p_head.text = heading
                p_head.font.bold = True
                p_head.font.size = Pt(12.5)
                p_head.font.color.rgb = RGBColor(255, 255, 255)
                p_head.space_after = Pt(2)
                
                p_body = tf.add_paragraph()
                p_body.text = body
                p_body.font.size = Pt(10.5)
                p_body.font.color.rgb = RGBColor(220, 220, 225)
                p_body.space_after = Pt(6)

prs.save(OUTPUT_PATH)
print(f"SUCCESS: Remade presentation saved to: {OUTPUT_PATH}")
